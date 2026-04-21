import io
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import AzureOpenAI, OpenAI
from app.config import settings

if settings.AZURE_OPENAI_EMBEDDING_ENDPOINT:
    client = AzureOpenAI(
        api_key=settings.AZURE_OPENAI_EMBEDDING_KEY or settings.OPENAI_API_KEY,
        api_version=settings.AZURE_OPENAI_API_VERSION,
        azure_endpoint=settings.AZURE_OPENAI_EMBEDDING_ENDPOINT
    )
    EMBEDDING_MODEL = settings.AZURE_OPENAI_EMBEDDING_DEPLOYMENT
else:
    client = OpenAI(api_key=settings.OPENAI_API_KEY)
    EMBEDDING_MODEL = "text-embedding-3-small"

# --- Extracción de texto ---

def extract_text(file_bytes: bytes, file_type: str) -> str:
    if file_type == "pdf":
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    elif file_type in ("docx",):
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)
    elif file_type == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    elif file_type in ("txt", "md"):
        return file_bytes.decode("utf-8", errors="ignore")
    else:
        raise ValueError(f"Tipo de archivo no soportado: {file_type}")

# --- Chunking ---

def chunk_text(text: str, chunk_size: int = 800, chunk_overlap: int = 150) -> list[str]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    return splitter.split_text(text)

# --- Embeddings ---

def get_embeddings(texts: list[str]) -> list[list[float]]:
    """Genera embeddings en batch (máximo 2048 textos por llamada)."""
    response = client.embeddings.create(
        model=EMBEDDING_MODEL,
        input=texts,
    )
    return [item.embedding for item in response.data]
